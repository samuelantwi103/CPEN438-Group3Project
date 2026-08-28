#!/usr/bin/env python3
"""
build_final_defence_deck.py — generates week4/presentation/final_defence_deck.pptx

Project 1: Roofline Reckoning, CPEN 438, Group 3. Week 4 deliverable
(brief Part V §15: "Presentation slides (paper-review deck AND final
defence deck) present in presentation/" — this is the SECOND, separate
deck; the paper-review deck from Week 1 stays where it is and is not
duplicated here).

Content pulled directly from week4/report/technical_report.md and the
real measured data cited there — not reinvented for the slides.

Run: python build_final_defence_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent / "final_defence_deck.pptx"

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0xF1, 0x8F, 0x01)
GREY = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title, bullets, presenter=None):
    slide = prs.slides.add_slide(BLANK)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.6))
    btf = body.text_frame
    btf.word_wrap = True
    for i, (level, text) in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = text
        para.level = level
        para.font.size = Pt(19 if level == 0 else 16)
        para.font.color.rgb = NAVY if level == 0 else GREY
        para.space_after = Pt(7)

    if presenter:
        pbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(6), Inches(0.4))
        ptf = pbox.text_frame
        ptf.text = f"Presenter: {presenter}"
        ptf.paragraphs[0].font.size = Pt(12)
        ptf.paragraphs[0].font.color.rgb = ACCENT
        ptf.paragraphs[0].font.italic = True
    return slide


# ---- Title ----
slide = prs.slides.add_slide(BLANK)
box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.8))
tf = box.text_frame
tf.word_wrap = True
tf.text = "Roofline Reckoning"
tf.paragraphs[0].font.size = Pt(40)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = NAVY
p2 = tf.add_paragraph()
p2.text = "Quantitative Performance Characterisation of Ghanaian Fintech and Climate Workloads"
p2.font.size = Pt(20)
p2.font.color.rgb = GREY
p3 = tf.add_paragraph()
p3.text = "Final Defence — Project 1, Group 3, CPEN 438"
p3.font.size = Pt(16)
p3.font.italic = True
p3.font.color.rgb = ACCENT
p4 = tf.add_paragraph()
p4.text = ("Samuel Kojo Anafi Antwi (11164744) · Amponsah Jonathan Boadu (11293871) · "
           "Kumi Kelvin Gyabaah (11012343) · Obed Ninson (11238291) · "
           "Theophilus Owusu-Manu (10985130)")
p4.font.size = Pt(14)
p4.font.color.rgb = GREY

add_slide("The Problem", [
    (0, "Hardware procurement decisions for Ghanaian fintech/public-sector systems are made by analogy to generic benchmarks (e.g. SPEC CPU2017)."),
    (0, "Generic benchmarks are compute-bound, portable, predictable-code-path programs by design — they do not resemble transaction matching, small-model regression, or irregular geospatial interpolation."),
    (0, "Claim to be tested: local, measured, workload-specific benchmarking tells a materially different story than a generic one."),
], presenter="Samuel Kojo Anafi Antwi")

add_slide("What We Built: GH-Bench", [
    (0, "Three kernels, each a stand-in for a real Ghanaian workload class:"),
    (1, "momo_match — mobile-money transaction matching (hash join)"),
    (1, "cocobod_yield_regression — cocoa yield prediction (OLS regression)"),
    (1, "rainfall_interpolate — Volta-basin rainfall gap-filling (IDW interpolation)"),
    (0, "Measured — not estimated — using real hardware performance counters via the Linux perf_event_open syscall."),
    (0, "Two hardware configurations: 2-core and 8-core, taskset-restricted on one physical machine."),
], presenter="Kumi Kelvin Gyabaah")

add_slide("Measurement Methodology", [
    (0, "Harness reads instructions, cycles, and cache-miss counters directly from CPU hardware — not a generic timer."),
    (0, "Timed region discipline: only the kernel call itself is measured, never I/O or data generation."),
    (0, "Real problem hit and solved: the fine-grained LLC-miss counter silently read 0 under our WSL2 virtualised environment — caught with a dedicated cache-busting test, substituted for a working generic counter."),
    (0, "Five repetitions per cell after one documented warm-up run; 36 total measured runs in the main matrix."),
], presenter="Amponsah Jonathan Boadu")

add_slide("Results — CPU-Time Validation", [
    (0, "All six (kernel x configuration) cells validated against an independent analytical prediction, within the required +/-15% tolerance."),
    (1, "Worst case: 9.41% (cocobod_yield_regression, configA)"),
    (1, "Best case: 0.68% (rainfall_interpolate, configA)"),
    (0, "This confirms the measurement pipeline is internally consistent — not yet independent proof the formula predicts anything (see manual worked example for that)."),
], presenter="Obed Ninson")

add_slide("Results — Roofline: All Three Kernels Are Memory-Bound", [
    (0, "Every kernel, at both configurations, sits below its measured compute ceiling — none reach the compute-bound region."),
    (1, "cocobod_yield_regression: 66.9-71.5% of ceiling (closest to compute-bound of the three)"),
    (1, "rainfall_interpolate: 15.5-36.4% of ceiling"),
    (1, "momo_match: 19.1-29.1% of ceiling (furthest — latency-bound, not just bandwidth-bound)"),
    (0, "This is exactly the finding a SPEC-style compute-bound benchmark would never surface for these workload classes."),
], presenter="Obed Ninson")

add_slide("Results — Amdahl's Law and the OpenMP Extension", [
    (0, "momo_match's serial (build) vs. parallel (probe) phases separately measured: 33.9% / 66.1% split."),
    (0, "Amdahl projects up to 2.37x speedup at 8x local parallelism."),
    (0, "MEASURED OpenMP speedup: 1.51x (2T), 1.77x (4T), 1.49x (8T) — sublinear, and actually WORSE at 8 threads than 4."),
    (0, "Explanation: random hash-table probes are memory-latency-bound; beyond 4 physical cores, hyperthreads compete for the same memory pipeline rather than adding independent throughput."),
], presenter="Amponsah Jonathan Boadu")

add_slide("Innovation: A Self-Calibrating Adaptive Kernel", [
    (0, "Problem: a FIXED thread-count decision is wrong at small n — overhead exceeds the work itself."),
    (0, "Our fix: measure, at runtime, this machine's actual thread-launch overhead AND this workload's actual per-transaction cost, then decide."),
    (0, "At our project's real dataset size (n=20,000): a naive always-parallel version would be 8.5x SLOWER. Our adaptive kernel avoids this automatically."),
    (0, "At n=500,000+: correctly switches to parallel, captures ~2x real speedup."),
    (0, "RAPL energy measurement (the brief's other suggested idea) was investigated and confirmed unavailable in our WSL2 environment before choosing this direction."),
], presenter="Kumi Kelvin Gyabaah")

add_slide("Honest Limitations", [
    (0, "One physical machine, not two — configs are taskset-restricted core counts, a brief-sanctioned but weaker substitution."),
    (0, "Cache-miss counter is a generic hardware event, not LLC-specific, due to a virtualisation limitation."),
    (0, "Measured (not vendor-datasheet) Roofline peaks — a more aggressive AVX2/FMA compute microbenchmark would likely raise the compute ceiling further."),
    (0, "Real run-to-run timing variance observed and documented throughout — exact figures are one run's output, not fixed constants."),
    (0, "MATLAB became unavailable on our machine during Week 4 — final validation numbers were reproduced directly in Python from the same formulas, disclosed explicitly."),
], presenter="Samuel Kojo Anafi Antwi")

add_slide("Conclusion", [
    (0, "GH-Bench demonstrates, with real measured data, that generic benchmark methodology misses workload classes relevant to Ghanaian fintech/climate computing — all three of our kernels are memory-bound, a finding SPEC's own compute-bound selection criteria would never surface."),
    (0, "CPU-time equation validated on all 6 cells; Amdahl correctly bounded the measured OpenMP result; the innovation kernel delivered a concrete, measured benefit motivated directly by that finding."),
    (0, "What we'd do with one more week: an AVX2/FMA compute microbenchmark to test whether the Roofline ceiling classification changes, and investigate rainfall_interpolate's unexplained cache-miss jump at 8 cores."),
], presenter="Samuel Kojo Anafi Antwi")

add_slide("Live Defence: Unseen Configuration", [
    (0, "[Placeholder — the brief specifies the instructor provides an unseen third hardware/software configuration parameter at defence time.]"),
    (0, "Team should be ready to: describe, live, how the harness/experiment-matrix script would be extended to accept it; predict qualitatively where each kernel would land on the Roofline plot before running it."),
], presenter="Whole team")

prs.save(str(OUT))
print(f"wrote {OUT} ({len(prs.slides)} slides)")
